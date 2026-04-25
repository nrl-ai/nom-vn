"""Generate synthetic Vietnamese Office documents for parser tests.

Idempotent. Re-running overwrites the fixtures. Run::

    python benchmarks/data/office_vi/_generate.py

Outputs (all under this directory):
- ``hop_dong.docx``    — DOCX with paragraphs + a 3-row table.
- ``so_sach.xlsx``     — XLSX with three sheets (each non-trivial).
- ``thuyet_trinh.pptx``— PPTX with three slides + speaker notes.
- ``ground_truth.json``— what the parser should extract from each.

License: synthetic content authored here, public domain. Use freely.
"""

from __future__ import annotations

import json
from pathlib import Path

DIR = Path(__file__).resolve().parent


def _make_docx() -> dict:
    from docx import Document

    doc = Document()
    doc.add_heading("Hợp đồng dịch vụ tư vấn pháp lý", 0)
    doc.add_paragraph("Số hợp đồng: HĐ-2026/045. Ngày ký: 14 tháng 3 năm 2026.")
    doc.add_paragraph(
        "Bên A: Công ty TNHH Pháp lý Hồng Hà — đại diện bởi ông Nguyễn Văn A, chức vụ Giám đốc."
    )
    doc.add_paragraph(
        "Bên B: Công ty Cổ phần Xây dựng Sông Hồng — đại diện bởi bà Lê Thị B, chức vụ Tổng giám đốc."
    )
    doc.add_heading("Điều 1 — Phạm vi công việc", level=1)
    doc.add_paragraph(
        "Bên A cam kết cung cấp dịch vụ tư vấn pháp lý cho Bên B trong các lĩnh vực: hợp đồng thương mại, "
        "tranh chấp dân sự, sở hữu trí tuệ, và pháp luật doanh nghiệp."
    )
    doc.add_heading("Điều 2 — Giá trị hợp đồng", level=1)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Light Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Hạng mục"
    hdr[1].text = "Đơn vị"
    hdr[2].text = "Giá trị (VND)"
    rows = [
        ("Tư vấn cố định hàng tháng", "tháng × 12", "180.000.000"),  # noqa: RUF001
        ("Tư vấn theo vụ việc", "ước tính", "120.000.000"),
        ("Tổng cộng", "", "300.000.000"),
    ]
    for i, (a, b, c) in enumerate(rows, start=1):
        cells = table.rows[i].cells
        cells[0].text = a
        cells[1].text = b
        cells[2].text = c
    doc.add_heading("Điều 3 — Phạt vi phạm", level=1)
    doc.add_paragraph(
        "Mức phạt vi phạm hợp đồng: 8% giá trị hợp đồng nếu vi phạm nghiêm trọng nghĩa vụ thanh toán."
    )
    out = DIR / "hop_dong.docx"
    doc.save(out)
    return {
        "file": out.name,
        "expected_paragraphs_contain": [
            "HĐ-2026/045",
            "Hồng Hà",
            "Sông Hồng",
            "300.000.000",
            "8%",
        ],
        "expected_min_pages": 7,  # paragraphs + table rows
    }


def _make_xlsx() -> dict:
    from openpyxl import Workbook

    wb = Workbook()
    # Sheet 1 — finance summary
    ws1 = wb.active
    ws1.title = "Doanh thu Q1"
    ws1.append(["Tháng", "Doanh thu (VND)", "Chi phí (VND)", "Lợi nhuận (VND)"])
    for label, rev, cost in [
        ("Tháng 1", 850_000_000, 620_000_000),
        ("Tháng 2", 920_000_000, 640_000_000),
        ("Tháng 3", 1_050_000_000, 700_000_000),
    ]:
        ws1.append([label, rev, cost, rev - cost])
    # Sheet 2 — staff
    ws2 = wb.create_sheet("Nhân sự")
    ws2.append(["Họ và tên", "Chức vụ", "Phòng ban"])
    for name, role, dept in [
        ("Trần Văn C", "Trưởng phòng", "Kỹ thuật"),
        ("Lê Thị D", "Nhân viên", "Kinh doanh"),
        ("Phạm Minh E", "Phó giám đốc", "Tài chính"),
    ]:
        ws2.append([name, role, dept])
    # Sheet 3 — projects
    ws3 = wb.create_sheet("Dự án")
    ws3.append(["Mã dự án", "Tên dự án", "Trạng thái"])
    for code, name, status in [
        ("DA-001", "Cầu vượt đường Trường Chinh", "Đang triển khai"),
        ("DA-002", "Toà nhà văn phòng Lê Lợi", "Hoàn thành"),
    ]:
        ws3.append([code, name, status])
    out = DIR / "so_sach.xlsx"
    wb.save(out)
    return {
        "file": out.name,
        "expected_sheets": ["Doanh thu Q1", "Nhân sự", "Dự án"],
        "expected_pages_contain": ["Doanh thu Q1", "Nhân sự", "Dự án"],
        "expected_text_contains": ["1050000000", "Trường Chinh", "Lê Thị D"],
    }


def _make_pptx() -> dict:
    from pptx import Presentation

    prs = Presentation()
    # Slide 1 — title only
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Báo cáo kết quả Quý 1 năm 2026"
    s1.placeholders[1].text = "Phòng Kế hoạch & Tài chính"
    # Accessing notes_slide creates it on demand — don't gate on has_notes_slide.
    s1.notes_slide.notes_text_frame.text = "Bắt đầu bằng tóm tắt 30 giây."
    # Slide 2 — bullet list (Title + Content layout)
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Doanh thu vượt kế hoạch 12%"
    body = s2.placeholders[1].text_frame
    body.text = "Doanh thu Q1: 2.820 tỷ VND"
    body.add_paragraph().text = "Lợi nhuận sau thuế: 760 tỷ VND"
    body.add_paragraph().text = "Tăng trưởng so với cùng kỳ năm trước: +18%"
    s2.notes_slide.notes_text_frame.text = "Nhấn mạnh tăng trưởng tháng 3."
    # Slide 3 — single-sentence body
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Kế hoạch Q2"
    s3.placeholders[1].text_frame.text = "Mở rộng thị trường miền Trung và đầu tư vào ba dự án mới."
    out = DIR / "thuyet_trinh.pptx"
    prs.save(out)
    return {
        "file": out.name,
        "expected_n_slides": 3,
        "expected_titles": [
            "Báo cáo kết quả Quý 1 năm 2026",
            "Doanh thu vượt kế hoạch 12%",
            "Kế hoạch Q2",
        ],
        "expected_notes_contain": ["30 giây", "tháng 3"],
    }


def main() -> None:
    DIR.mkdir(parents=True, exist_ok=True)
    truth = {
        "docx": _make_docx(),
        "xlsx": _make_xlsx(),
        "pptx": _make_pptx(),
    }
    (DIR / "ground_truth.json").write_text(
        json.dumps(truth, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    print(f"wrote fixtures to {DIR}")


if __name__ == "__main__":
    main()
