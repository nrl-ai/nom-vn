"""Pipeline stage implementations.

All six default stages are real as of v0.0.3:

  - Load (pure stdlib, magic-byte format detection)
  - Parse (pdfplumber MIT default; pymupdf AGPL opt-in)
  - OCR (pytesseract + vie traineddata)
  - Normalize (nom.text)
  - Extract (LLM via nom.llm + Pydantic schema with retries)
  - Validate (nom.doc.schemas SchemaResolver, Pydantic v2)

OSS prior art studied while designing this module:

- **LangChain ``DocumentLoader`` family** — same Load/Parse split we use
  here. We adopt the per-source-format dispatch but keep our Stage
  Protocol orthogonal to LangChain's runnable abstraction.
- **LlamaIndex ``IngestionPipeline``** — composable transformations on
  a Document. Our ``Pipeline`` mirrors that shape (stages thread a
  ``Context`` rather than transforming a single object).
- **Unstructured.io** — ``partition_pdf`` / ``partition_image`` for
  format-aware extraction. Our magic-byte detection is influenced by
  Unstructured's auto-detection logic but uses pure stdlib.
- **Instructor** — the retry-with-error-feedback pattern in the Extract
  stage is lifted from this library, reimplemented in ~30 LOC so we
  don't add the dep.

See ``docs/pipeline.md`` for the per-stage component picks and rationale.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from nom.doc.pipeline import Context

__all__ = [
    "OCR",
    "Extract",
    "Load",
    "Normalize",
    "Parse",
    "Validate",
]


# Canonical mime/format detection. Magic-bytes-aware for the formats we care
# about; falls back to extension. No third-party deps.
_PDF_MAGIC = b"%PDF-"
_IMAGE_MAGICS: dict[bytes, str] = {
    b"\x89PNG\r\n\x1a\n": "image/png",
    b"\xff\xd8\xff": "image/jpeg",  # JPEG magic prefix (FF D8 FF)
    b"GIF87a": "image/gif",
    b"GIF89a": "image/gif",
    b"II*\x00": "image/tiff",  # little-endian TIFF
    b"MM\x00*": "image/tiff",  # big-endian TIFF
    b"RIFF": "image/webp",  # WEBP starts with RIFF (full check below)
    b"BM": "image/bmp",
}

_TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv"}
_HTML_EXTENSIONS = {".html", ".htm", ".xhtml"}
_JSON_EXTENSIONS = {".json", ".jsonl", ".ndjson"}

# Office Open XML magic prefix is ``PK\x03\x04`` — the same as any ZIP.
# We disambiguate by extension when the file is path-based, and by
# inspecting the central directory when it's bytes (looking for the
# format-specific top-level entries).
_ZIP_MAGIC = b"PK\x03\x04"


def _detect_format(source: str | Path | bytes) -> str:
    """Return one of ``"pdf"``, ``"image"``, ``"docx"``, ``"xlsx"``,
    ``"pptx"``, ``"html"``, ``"json"``, ``"text"``.

    For bytes: magic-byte sniff, with ZIP-based formats disambiguated by
    inspecting the embedded central directory.
    For path: magic first, fall back to extension.
    """
    if isinstance(source, bytes):
        head = source[:16]
        fmt = _format_from_head(head, ext="")
        if fmt == "zip":
            return _format_from_zip_bytes(source)
        return fmt
    path = Path(source)
    if path.is_file():
        with path.open("rb") as f:
            head = f.read(16)
        ext = path.suffix.lower()
        fmt = _format_from_head(head, ext=ext)
        if fmt == "zip":
            # Path-based: trust the extension first; fall back to inspection.
            ext_fmt = _office_format_from_ext(ext)
            return ext_fmt or _format_from_zip_path(path)
        return fmt
    # Path doesn't exist (yet) — extension only
    return _format_from_extension(path.suffix.lower())


def _format_from_head(head: bytes, *, ext: str) -> str:
    if head.startswith(_PDF_MAGIC):
        return "pdf"
    for magic in _IMAGE_MAGICS:
        if head.startswith(magic):
            # WEBP requires the WEBP token at offset 8
            if magic == b"RIFF" and len(head) >= 12 and head[8:12] != b"WEBP":
                continue
            return "image"
    if head.startswith(_ZIP_MAGIC):
        return "zip"
    return _format_from_extension(ext)


def _format_from_extension(ext: str) -> str:
    if ext == ".pdf":
        return "pdf"
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".tiff", ".tif", ".bmp", ".webp"}:
        return "image"
    if ext in {".docx", ".docm"}:
        return "docx"
    if ext in {".xlsx", ".xlsm"}:
        return "xlsx"
    if ext in {".pptx", ".pptm"}:
        return "pptx"
    if ext in _HTML_EXTENSIONS:
        return "html"
    if ext in _JSON_EXTENSIONS:
        return "json"
    if ext in _TEXT_EXTENSIONS:
        return "text"
    # Default: treat unknowns as text. Keeps the pipeline running rather than
    # blowing up on a misnamed file.
    return "text"


def _office_format_from_ext(ext: str) -> str | None:
    """Return docx/xlsx/pptx for known Office extensions, else None."""
    return {
        ".docx": "docx",
        ".docm": "docx",
        ".xlsx": "xlsx",
        ".xlsm": "xlsx",
        ".pptx": "pptx",
        ".pptm": "pptx",
    }.get(ext)


def _format_from_zip_bytes(data: bytes) -> str:
    """Look inside a ZIP to detect Office Open XML variants.

    Word: ``word/document.xml``. Excel: ``xl/workbook.xml``.
    PowerPoint: ``ppt/presentation.xml``.
    """
    import io
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = set(zf.namelist())
    except zipfile.BadZipFile:
        return "text"
    if "word/document.xml" in names:
        return "docx"
    if "xl/workbook.xml" in names:
        return "xlsx"
    if "ppt/presentation.xml" in names:
        return "pptx"
    return "text"


def _format_from_zip_path(path: Path) -> str:
    return _format_from_zip_bytes(path.read_bytes())


class Load:
    """Detect input format and route to the next stage.

    v0.0.3 implementation: pure stdlib. Reads ``ctx.source`` (path/bytes),
    sets ``ctx.fmt`` ∈ {``"pdf"``, ``"image"``, ``"text"``}, populates
    ``ctx.metadata["bytes"]`` if source was a path so downstream stages can
    read it without a second disk hit.
    """

    name = "Load"

    def run(self, ctx: Context) -> Context:
        fmt = _detect_format(ctx.source)
        ctx.fmt = fmt
        # Cache bytes on context if we read from disk; downstream stages reuse.
        if isinstance(ctx.source, str | Path):
            path = Path(ctx.source)
            if path.is_file():
                ctx.metadata["path"] = path
                ctx.metadata["bytes"] = path.read_bytes()
        elif isinstance(ctx.source, bytes):
            ctx.metadata["bytes"] = ctx.source
        return ctx


class Parse:
    """Extract native text + layout from PDFs (skip if image/text).

    v0.0.3 implementation: uses ``pdfplumber`` (MIT) when installed.
    pdfplumber is slower than PyMuPDF (AGPL) but our default keeps us in
    permissive-license territory. Power users can swap to PyMuPDF.

    For ``"image"`` inputs: leaves ``ctx.pages_text`` empty and adds a single
    OCR-needed entry — the OCR stage handles it.

    For ``"text"`` inputs: reads ``ctx.metadata["bytes"]`` as UTF-8 and
    populates ``ctx.pages_text`` with one entry.

    Args:
        backend: ``"pdfplumber"`` (default, MIT) or ``"pymupdf"`` (AGPL,
            faster — opt-in for users who can comply with the license).
    """

    name = "Parse"

    def __init__(self, backend: str = "pdfplumber") -> None:
        if backend not in ("pdfplumber", "pymupdf"):
            raise ValueError(f"backend must be 'pdfplumber' or 'pymupdf', got {backend!r}")
        self.backend = backend

    def run(self, ctx: Context) -> Context:
        if ctx.fmt is None:
            raise RuntimeError("Parse stage requires ctx.fmt to be set — run Load first.")

        if ctx.fmt == "text":
            data = ctx.metadata.get("bytes", b"")
            ctx.pages_text = [data.decode("utf-8", errors="replace")]
            ctx.text = ctx.pages_text[0]
            return ctx

        if ctx.fmt == "image":
            # Parse skips images — OCR stage handles them.
            ctx.pages_text = [""]
            ctx.needs_ocr = [0]
            return ctx

        if ctx.fmt == "html":
            self._parse_html(ctx)
            return ctx
        if ctx.fmt == "json":
            self._parse_json(ctx)
            return ctx
        if ctx.fmt == "docx":
            self._parse_docx(ctx)
            return ctx
        if ctx.fmt == "xlsx":
            self._parse_xlsx(ctx)
            return ctx
        if ctx.fmt == "pptx":
            self._parse_pptx(ctx)
            return ctx

        # PDF path
        if self.backend == "pdfplumber":
            self._parse_pdf_pdfplumber(ctx)
        else:
            self._parse_pdf_pymupdf(ctx)
        return ctx

    # ------------------------------------------------------------------
    # Office Open XML parsers — pure-Python deps, MIT-licensed
    # ------------------------------------------------------------------

    def _parse_docx(self, ctx: Context) -> None:
        try:
            from docx import Document  # python-docx
        except ImportError as exc:
            raise ImportError(
                "Parse for .docx requires python-docx. Install with: pip install nom-vn[doc]"
            ) from exc
        import io

        data = ctx.metadata.get("bytes")
        if data is None:
            raise RuntimeError("Parse: no bytes available on context.")
        doc = Document(io.BytesIO(data))
        # Paragraphs first; then any tables (each cell as text).
        parts: list[str] = [p.text for p in doc.paragraphs if p.text]
        for tbl in doc.tables:
            for row in tbl.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells)
                if row_text.strip(" |"):
                    parts.append(row_text)
        ctx.pages_text = parts
        ctx.text = "\n\n".join(parts)

    def _parse_xlsx(self, ctx: Context) -> None:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ImportError(
                "Parse for .xlsx requires openpyxl. Install with: pip install nom-vn[doc]"
            ) from exc
        import io

        data = ctx.metadata.get("bytes")
        if data is None:
            raise RuntimeError("Parse: no bytes available on context.")
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        # One "page" per sheet; each sheet rendered as TSV-ish lines.
        # data_only=True follows formula results, not the formula text.
        pages: list[str] = []
        for sheet in wb.worksheets:
            lines: list[str] = [f"# {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    lines.append("\t".join(cells))
            pages.append("\n".join(lines))
        wb.close()
        ctx.pages_text = pages
        ctx.text = "\n\n".join(pages)

    def _parse_pptx(self, ctx: Context) -> None:
        try:
            from pptx import Presentation  # python-pptx
        except ImportError as exc:
            raise ImportError(
                "Parse for .pptx requires python-pptx. Install with: pip install nom-vn[doc]"
            ) from exc
        import io

        data = ctx.metadata.get("bytes")
        if data is None:
            raise RuntimeError("Parse: no bytes available on context.")
        prs = Presentation(io.BytesIO(data))
        # One page per slide. Format per slide:
        #   <title>          (from slide.shapes.title; may be empty)
        #   <body line 1>    (from non-title text frames)
        #   <body line 2>
        #   _notes: <notes>  (only if speaker notes present)
        # We don't prepend a "Slide N" marker because viewers know the
        # slide index from the page's position in the list.
        pages: list[str] = []
        for slide in prs.slides:
            title = ""
            try:
                if slide.shapes.title is not None:
                    title = (slide.shapes.title.text or "").strip()
            except Exception:
                pass
            body_lines: list[str] = []
            for shape in slide.shapes:
                tf = getattr(shape, "text_frame", None)
                if tf is None or not tf.text.strip():
                    continue
                # Skip the title shape — already captured above.
                if shape == slide.shapes.title:
                    continue
                body_lines.append(tf.text.strip())
            notes = ""
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            parts: list[str] = []
            if title:
                parts.append(title)
            parts.extend(body_lines)
            if notes:
                parts.append(f"_notes: {notes}")
            pages.append("\n".join(parts) if parts else "(empty slide)")
        ctx.pages_text = pages
        ctx.text = "\n\n".join(pages)

    # ------------------------------------------------------------------
    # HTML / JSON — stdlib parsers, no extra deps
    # ------------------------------------------------------------------

    def _parse_html(self, ctx: Context) -> None:
        from html.parser import HTMLParser

        data = ctx.metadata.get("bytes", b"")
        text = data.decode("utf-8", errors="replace")

        # Minimal HTML stripper — drop <script>/<style>, collect text.
        # No third-party dep; XSS-safe (we never re-emit HTML).
        class _Stripper(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self._skip = 0
                self.parts: list[str] = []

            def handle_starttag(self, tag: str, _attrs: list[tuple[str, str | None]]) -> None:
                if tag in ("script", "style", "noscript"):
                    self._skip += 1

            def handle_endtag(self, tag: str) -> None:
                if tag in ("script", "style", "noscript") and self._skip > 0:
                    self._skip -= 1

            def handle_data(self, data: str) -> None:
                if self._skip == 0:
                    self.parts.append(data)

        s = _Stripper()
        s.feed(text)
        joined = "\n".join(p.strip() for p in s.parts if p.strip())
        ctx.pages_text = [joined]
        ctx.text = joined

    def _parse_json(self, ctx: Context) -> None:
        """Render JSON / JSONL as readable lines.

        Each top-level object → one block. JSONL is line-delimited so
        each line becomes one entry. Strings are kept as-is; nested
        objects are JSON-pretty-printed so structure is searchable.
        """
        import json as _json

        data = ctx.metadata.get("bytes", b"").decode("utf-8", errors="replace")
        records: list[Any] = []
        # Try JSONL first.
        looks_jsonl = "\n" in data.strip() and data.strip().startswith(("{", "["))
        if looks_jsonl:
            for line in data.splitlines():
                line = line.strip()
                if not line:
                    continue
                try:
                    records.append(_json.loads(line))
                except _json.JSONDecodeError:
                    records.append(line)
        else:
            try:
                obj = _json.loads(data)
                records = obj if isinstance(obj, list) else [obj]
            except _json.JSONDecodeError:
                records = [data]
        pages: list[str] = []
        for r in records:
            if isinstance(r, str):
                pages.append(r)
            else:
                pages.append(_json.dumps(r, ensure_ascii=False, indent=2))
        ctx.pages_text = pages
        ctx.text = "\n\n".join(pages)

    def _parse_pdf_pdfplumber(self, ctx: Context) -> None:
        try:
            import pdfplumber
        except ImportError as exc:
            raise ImportError(
                "Parse(backend='pdfplumber') requires pdfplumber. "
                "Install with: pip install nom-vn[doc]"
            ) from exc

        data = ctx.metadata.get("bytes")
        if data is None:
            raise RuntimeError("Parse: no bytes available on context.")

        import io

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages: list[str] = []
            layouts: list[dict[str, Any]] = []
            needs_ocr: list[int] = []
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                pages.append(text)
                layouts.append(
                    {
                        "page": i,
                        "width": page.width,
                        "height": page.height,
                        "n_chars": len(text),
                    }
                )
                if not text.strip():
                    needs_ocr.append(i)

        ctx.pages_text = pages
        ctx.pages_layout = layouts
        ctx.needs_ocr = needs_ocr
        # Concatenated text for stages that don't care about page boundaries.
        ctx.text = "\n\n".join(pages)

    def _parse_pdf_pymupdf(self, ctx: Context) -> None:
        try:
            import fitz
        except ImportError as exc:
            raise ImportError(
                "Parse(backend='pymupdf') requires PyMuPDF (AGPL). "
                "Install with: pip install pymupdf"
            ) from exc

        data = ctx.metadata.get("bytes")
        if data is None:
            raise RuntimeError("Parse: no bytes available on context.")

        doc = fitz.open(stream=data, filetype="pdf")
        try:
            pages: list[str] = []
            layouts: list[dict[str, Any]] = []
            needs_ocr: list[int] = []
            for i, page in enumerate(doc):
                text = page.get_text()
                pages.append(text)
                layouts.append(
                    {
                        "page": i,
                        "width": page.rect.width,
                        "height": page.rect.height,
                        "n_chars": len(text),
                    }
                )
                if not text.strip():
                    needs_ocr.append(i)
        finally:
            doc.close()

        ctx.pages_text = pages
        ctx.pages_layout = layouts
        ctx.needs_ocr = needs_ocr
        ctx.text = "\n\n".join(pages)


class OCR:
    """Run OCR on pages flagged as scans (or on direct image inputs).

    v0.0.3 implementation: Tesseract via ``pytesseract`` (Apache 2.0).
    Tesseract is the most-audited OCR engine in the open-source world,
    and pytesseract is a thin wrapper (~hundreds of lines). Tesseract
    binary is system-installed: ``apt install tesseract-ocr tesseract-ocr-vie``
    on Debian/Ubuntu, or ``brew install tesseract tesseract-lang`` on
    macOS.

    For Vietnamese, this stage runs Tesseract with ``-l vie``. Per
    upstream docs, the ``vie.traineddata`` was trained on Times New Roman,
    Arial, Verdana, Courier New — accuracy is best on those fonts and
    inputs scanned at 200-400 DPI.

    Future v0.x backends in docs/pipeline.md (VietOCR Transformer,
    PaddleOCR PP-OCRv5) will integrate as opt-in alternatives once we
    have measured accuracy comparisons on a curated VN scan corpus.

    Reads:
        - ``ctx.fmt`` — must be ``"image"`` (else this is a no-op).
        - ``ctx.metadata["bytes"]`` — raw image bytes from Load.
        - Or for PDF inputs: pages flagged in ``ctx.needs_ocr`` (v0.1.1).

    Writes:
        - ``ctx.pages_text`` — the OCR'd text per page.
        - ``ctx.text`` — concatenated.
        - ``ctx.needs_ocr`` — cleared.

    Args:
        lang: Tesseract language code. Default ``"vie"``. Use ``"vie+eng"``
            for mixed Vietnamese/English documents.
        config: extra Tesseract config flags (e.g. ``"--psm 6"`` for
            uniform block of text). Default empty.
    """

    name = "OCR"

    def __init__(self, lang: str = "vie", config: str = "") -> None:
        self.lang = lang
        self.config = config

    def run(self, ctx: Context) -> Context:
        if ctx.fmt is None:
            raise RuntimeError("OCR requires ctx.fmt — run Load first.")

        # No-op for non-image, non-OCR-flagged inputs.
        if ctx.fmt != "image" and not ctx.needs_ocr:
            return ctx

        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            missing = "pytesseract" if "pytesseract" in str(exc) else "Pillow"
            raise ImportError(
                f"OCR stage requires {missing}. Install with: pip install nom-vn[doc]"
            ) from exc

        if ctx.fmt == "image":
            data = ctx.metadata.get("bytes")
            if data is None:
                raise RuntimeError("OCR: no image bytes available on context.")
            text = self._ocr_bytes(data, pytesseract, image_module=Image)
            ctx.pages_text = [text]
            ctx.text = text
            ctx.needs_ocr = []
        elif ctx.needs_ocr:
            # PDF pages with no text layer — rasterize each flagged page via
            # pdfplumber and OCR the rendered image. 200 DPI is the standard
            # speed/accuracy compromise for Tesseract (300 DPI is its docs'
            # ideal, but 200 keeps a 5-page demo under a second).
            data = ctx.metadata.get("bytes")
            if data is None:
                raise RuntimeError("OCR: no PDF bytes available on context.")
            import io

            import pdfplumber

            pages = list(ctx.pages_text)
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                for page_idx in ctx.needs_ocr:
                    page_img = pdf.pages[page_idx].to_image(resolution=200).original
                    if page_img.mode not in ("RGB", "L"):
                        page_img = page_img.convert("RGB")
                    pages[page_idx] = pytesseract.image_to_string(
                        page_img, lang=self.lang, config=self.config
                    ).strip()
            ctx.pages_text = pages
            ctx.text = "\n\n".join(pages)
            ctx.needs_ocr = []

        return ctx

    def _ocr_bytes(self, data: bytes, pytesseract: Any, *, image_module: Any) -> str:
        import io

        image = image_module.open(io.BytesIO(data))
        # Tesseract works best on RGB/grayscale; convert if needed.
        if image.mode not in ("RGB", "L"):
            image = image.convert("RGB")
        text: str = pytesseract.image_to_string(
            image,
            lang=self.lang,
            config=self.config,
        )
        # Tesseract sometimes emits trailing form-feeds and excessive
        # whitespace — strip and collapse.
        return text.strip()


class Normalize:
    """Apply ``nom.text`` cleanup to the parsed text.

    What it does (v0.0.3 implementation):
      1. Joins ``ctx.pages_text`` into ``ctx.text`` (page-separator: blank line).
      2. Applies Unicode NFC normalization.
      3. Applies VN-aware text normalization (whitespace + punctuation cleanup).
      4. Optionally applies diacritic restoration (off by default — the rule-
         based path corrupts text with high diacritic content; useful only on
         OCR-stripped text).

    Args:
        restore_diacritics: when True, apply ``nom.text.fix_diacritics`` to
            the joined text. Default False because the v0.0.1 rule-based
            implementation only helps OCR-stripped input — it can damage
            already-correct VN text. v0.0.3 plans an ML-backed restoration
            that will be safe to run unconditionally; until then this stays
            opt-in.
    """

    name = "Normalize"

    def __init__(self, restore_diacritics: bool = False) -> None:
        self.restore_diacritics = restore_diacritics

    def run(self, ctx: Context) -> Context:
        # Lazy import to keep the placeholder stages dep-free.
        from nom.text import fix_diacritics, normalize, text_normalize

        def _clean(s: str) -> str:
            out = text_normalize(normalize(s))
            return fix_diacritics(out) if self.restore_diacritics else out

        # Apply normalization per-page so we can preserve page boundaries in
        # ctx.text. text_normalize collapses internal whitespace, so joining
        # raw pages first would lose the page break.
        if ctx.pages_text:
            ctx.pages_text = [_clean(p) for p in ctx.pages_text]
            ctx.text = "\n\n".join(ctx.pages_text)
        elif ctx.text:
            # Defensive: if Parse didn't populate pages_text but did set
            # ctx.text directly (unusual), still normalize it.
            ctx.text = _clean(ctx.text)
        return ctx


class Extract:
    """Schema-driven LLM extraction with auto-retry on validation failure.

    v0.0.3 real implementation: builds a JSON schema from ``ctx.schema``,
    prompts the LLM with the parsed text + the schema, parses the JSON
    response, and retries with error feedback if the model produces
    invalid output. This is the same pattern as the ``instructor`` library
    (~30 LOC, no extra dep).

    Reads:
        - ``ctx.text`` — cleaned text from Normalize
        - ``ctx.schema`` — user-provided schema dict
    Writes:
        - ``ctx.output`` — raw dict from the LLM (Validate stage will
          coerce it through Pydantic).

    Args:
        llm: an ``LLM`` adapter from ``nom.llm`` (Ollama / OpenAI / Anthropic).
            Must have a ``.complete(prompt, schema=..., max_tokens=...)``
            method.
        max_retries: how many times to retry with error feedback when the
            LLM produces invalid JSON. Default 3.
        max_tokens: hint forwarded to the LLM. Default 2048.
    """

    name = "Extract"

    def __init__(
        self,
        llm: Any,
        *,
        max_retries: int = 3,
        max_tokens: int = 2048,
    ) -> None:
        if llm is None:
            raise ValueError(
                "Extract requires an LLM adapter. "
                "Use nom.llm.Ollama() for local inference, "
                "or implement the LLM protocol yourself."
            )
        if max_retries < 1:
            raise ValueError("max_retries must be >= 1")
        self.llm = llm
        self.max_retries = max_retries
        self.max_tokens = max_tokens

    def run(self, ctx: Context) -> Context:
        if not ctx.schema:
            raise RuntimeError(
                "Extract requires ctx.schema to be set. Provide one when calling Pipeline.run(...)."
            )
        if not ctx.text:
            raise RuntimeError("Extract requires ctx.text — did Parse and Normalize run?")

        from nom.doc.schemas import SchemaResolver

        resolver = SchemaResolver(ctx.schema)
        json_schema = resolver.json_schema()

        prompt = self._build_prompt(ctx.text, json_schema)
        last_error: str | None = None

        for attempt in range(1, self.max_retries + 1):
            response = self.llm.complete(
                prompt,
                schema=json_schema,
                max_tokens=self.max_tokens,
            )
            try:
                parsed = self._parse_json(response)
                ctx.output = parsed
                ctx.metadata["extract_attempts"] = attempt
                return ctx
            except (json.JSONDecodeError, ValueError) as exc:
                last_error = str(exc)
                # On retry, append error feedback so the model self-corrects.
                # This is the "instructor pattern" — see python.useinstructor.com.
                prompt = (
                    f"{prompt}\n\n"
                    f"=== Previous attempt produced invalid output ===\n"
                    f"Output snippet: {response[:200]}\n"
                    f"Error: {last_error}\n"
                    f"Please respond with ONLY a valid JSON object that "
                    f"strictly matches the schema. No prose, no markdown fences."
                )

        raise RuntimeError(
            f"Extract failed after {self.max_retries} attempts. Last error: {last_error}"
        )

    @staticmethod
    def _build_prompt(text: str, json_schema: dict[str, Any]) -> str:
        schema_str = json.dumps(json_schema, ensure_ascii=False, indent=2)
        return (
            "You are a Vietnamese document-extraction assistant. Extract the "
            "fields described by the JSON schema below from the document "
            "text. Respond with ONLY a valid JSON object — no prose, no "
            "markdown fences, no explanations.\n\n"
            "Vietnamese conventions to preserve:\n"
            "- Dates: 'ngày 14 tháng 3 năm 2025' or '14/3/2025' formats are OK.\n"
            "- Amounts: VND amounts use '.' as thousands separator (e.g., "
            "'1.500.000.000').\n"
            "- Diacritics: keep tone marks exactly as in the source.\n\n"
            f"=== Schema ===\n{schema_str}\n\n"
            f"=== Document ===\n{text}\n\n"
            "=== Response (JSON only) ==="
        )

    @staticmethod
    def _parse_json(response: str) -> dict[str, Any]:
        """Parse the LLM's JSON response, tolerant of markdown fences."""
        s = response.strip()
        # Strip ```json ... ``` and ``` ... ``` markdown fences if the model
        # ignored our instructions.
        if s.startswith("```"):
            # Drop the first fence line and the closing fence.
            lines = s.split("\n")
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines and lines[-1].strip() == "```":
                lines = lines[:-1]
            s = "\n".join(lines)
        result = json.loads(s)
        if not isinstance(result, dict):
            raise ValueError(f"Expected JSON object, got {type(result).__name__}")
        return result


class Validate:
    """Validate the LLM's structured output against the user-provided schema.

    v0.0.3 real implementation: uses :class:`nom.doc.schemas.SchemaResolver`
    to build a runtime Pydantic model from the user's schema dict, then
    validates ``ctx.output`` against it.

    Reads:
        - ``ctx.output`` — raw dict from Extract stage
        - ``ctx.schema`` — user-provided schema spec
    Writes:
        - ``ctx.output`` — validated + coerced dict (same keys, parsed values)
    """

    name = "Validate"

    def run(self, ctx: Context) -> Context:
        if not ctx.schema:
            raise RuntimeError(
                "Validate stage requires ctx.schema to be set. "
                "Provide a schema when calling Pipeline.run(...)."
            )
        if not ctx.output:
            raise RuntimeError(
                "Validate stage requires ctx.output to be populated by Extract. "
                "Did the Extract stage run?"
            )

        from nom.doc.schemas import SchemaResolver

        resolver = SchemaResolver(ctx.schema)
        ctx.output = resolver.validate(ctx.output)
        return ctx
