"""Format-preserving ``.pptx`` translation walker.

Walks every slide's shapes via ``python-pptx``. For each text frame
(title placeholders, body placeholders, text boxes, table cells, smart-
art labels), runs the same paragraph-level run-redistribution as the
docx walker: concatenate runs in a paragraph, translate as one unit,
write into the first run, clear the rest.

Trade-offs:

- **Mixed-style paragraphs** lose intra-paragraph styling. Same v0
  trade-off as the docx walker; deferred to a future word-alignment
  redistribution.
- **Speaker notes** (``slide.notes_slide``) are translated.
- **Charts and SmartArt** are out of scope for v0 — they live in
  separate XML streams that python-pptx doesn't fully expose.
- **Master / layout slides** are not visited — they only show through
  to the final render via inheritance, and translating them risks
  changing every slide's headers/footers in ways the user can't see.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING, Any

from nom.translate.base import Translator

if TYPE_CHECKING:
    pass

__all__ = ["PptxTranslationStats", "translate_pptx"]


@dataclass(frozen=True, slots=True)
class PptxTranslationStats:
    """Summary of a single :func:`translate_pptx` run."""

    paragraphs_translated: int
    paragraphs_skipped: int
    paragraphs_failed: int
    chars_in: int
    chars_out: int


def translate_pptx(
    src: Path | str,
    dst: Path | str,
    translator: Translator,
) -> PptxTranslationStats:
    """Translate text content of a ``.pptx``. Source untouched; ``dst``
    written fresh."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError(
            "Translating .pptx requires python-pptx. Install with: pip install nom-vn[doc]"
        ) from exc

    src_path = Path(src)
    dst_path = Path(dst)
    if not src_path.exists():
        raise FileNotFoundError(f"pptx source not found: {src_path}")

    prs = Presentation(str(src_path))
    counts = _Counts()

    for slide in prs.slides:
        for shape in slide.shapes:
            _translate_shape(shape, translator, counts)
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            _translate_text_frame(notes_tf, translator, counts)

    dst_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst_path))
    return counts.freeze()


@dataclass
class _Counts:
    translated: int = 0
    skipped: int = 0
    failed: int = 0
    chars_in: int = 0
    chars_out: int = 0

    def freeze(self) -> PptxTranslationStats:
        return PptxTranslationStats(
            paragraphs_translated=self.translated,
            paragraphs_skipped=self.skipped,
            paragraphs_failed=self.failed,
            chars_in=self.chars_in,
            chars_out=self.chars_out,
        )


def _translate_shape(shape: Any, translator: Translator, counts: _Counts) -> None:
    if shape.has_text_frame:
        _translate_text_frame(shape.text_frame, translator, counts)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _translate_text_frame(cell.text_frame, translator, counts)
    # Recurse into grouped shapes
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub in shape.shapes:
            _translate_shape(sub, translator, counts)


def _translate_text_frame(text_frame: Any, translator: Translator, counts: _Counts) -> None:
    for paragraph in text_frame.paragraphs:
        runs = list(paragraph.runs)
        source = "".join(r.text for r in runs) if runs else ""
        if not source:
            continue
        if not source.strip():
            counts.skipped += 1
            continue
        counts.chars_in += len(source)
        try:
            translated = translator.translate(source)
        except Exception:
            counts.failed += 1
            continue
        counts.chars_out += len(translated)
        runs[0].text = translated
        for run in runs[1:]:
            run.text = ""
        counts.translated += 1
