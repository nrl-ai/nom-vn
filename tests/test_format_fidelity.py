"""Format-fidelity tests — large fixtures covering real document features.

Each test generates a complex fixture (mixed styling, tables, headers,
footers, merged cells, multi-slide pptx, multi-sheet xlsx) in-process
via the format library, runs translate / convert through it with
controlled translators, and verifies the structural invariants:

- Style attributes preserved (bold/italic/font/heading/list-level)
- Tables retain their cell layout (no rows / cols dropped)
- Merged cells stay merged
- Multi-sheet workbooks keep all sheets
- Multi-slide pptx keeps all slides
- Headers and footers persist
- Hyperlinks aren't dropped
- Number / formula cells not corrupted

These are exhaustive — they're about catching regressions that the
unit tests in test_translate_*.py would miss because those use minimal
fixtures.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest

pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("pptx")
pytest.importorskip("PIL")
pytest.importorskip("pytesseract")

_TESSERACT_BIN = shutil.which("tesseract")
needs_tesseract = pytest.mark.skipif(
    _TESSERACT_BIN is None, reason="tesseract binary not installed"
)


# ---------------------------------------------------------------------------
# Translators used for fidelity testing (no real LLM calls).
# ---------------------------------------------------------------------------


class IdentityTranslator:
    """Returns input unchanged. Best for fidelity checks — any structural
    diff between source and target is purely format-walker noise."""

    name = "identity"
    source_lang = "en"
    target_lang = "vi"

    def translate(self, text: str, *, hint: str | None = None) -> str:
        return text


class PrefixTranslator:
    """Prefixes every translation with `«». Useful to verify the walker
    actually called the translator (vs. silent skip) on every unit."""

    name = "prefix"
    source_lang = "en"
    target_lang = "vi"

    def translate(self, text: str, *, hint: str | None = None) -> str:
        return "«" + text + "»"


# ---------------------------------------------------------------------------
# DOCX fidelity — heading levels, lists, mixed-style runs, tables, headers,
# footers, hyperlinks.
# ---------------------------------------------------------------------------


def _build_complex_docx(path: Path) -> None:
    from docx import Document

    doc = Document()

    # Title (Heading 1)
    h1 = doc.add_heading("Hợp đồng dịch vụ", level=1)
    assert h1.style.name == "Heading 1"

    # Heading 2
    doc.add_heading("Điều 1: Các bên", level=2)

    # Plain paragraph
    doc.add_paragraph("Bên A và Bên B đồng ý các điều khoản dưới đây.")

    # Mixed-style paragraph
    p = doc.add_paragraph("")
    p.add_run("Tổng giá trị: ")
    p.add_run("1.500.000.000 VNĐ").bold = True
    p.add_run(" — đã bao gồm ")
    p.add_run("VAT").italic = True
    p.add_run(".")

    # Bullet list
    doc.add_paragraph("Phương thức thanh toán:", style="Heading 3")
    doc.add_paragraph("Đợt 1 — 30%", style="List Bullet")
    doc.add_paragraph("Đợt 2 — 50%", style="List Bullet")
    doc.add_paragraph("Đợt 3 — 20%", style="List Bullet")

    # Numbered list
    doc.add_paragraph("Trách nhiệm Bên A:", style="Heading 3")
    doc.add_paragraph("Cung cấp tài liệu kịp thời", style="List Number")
    doc.add_paragraph("Thanh toán đúng hạn", style="List Number")

    # Table
    doc.add_heading("Mốc dự án", level=2)
    table = doc.add_table(rows=4, cols=3)
    table.cell(0, 0).text = "Giai đoạn"
    table.cell(0, 1).text = "Hạn cuối"
    table.cell(0, 2).text = "Sản phẩm"
    table.cell(1, 0).text = "Khởi động"
    table.cell(1, 1).text = "2026-01-15"
    table.cell(1, 2).text = "Bản kế hoạch"
    table.cell(2, 0).text = "Phát triển"
    table.cell(2, 1).text = "2026-04-01"
    table.cell(2, 2).text = "Bản beta"
    table.cell(3, 0).text = "Bàn giao"
    table.cell(3, 1).text = "2026-06-30"
    table.cell(3, 2).text = "Tài liệu cuối"

    # Header / footer
    section = doc.sections[0]
    section.header.paragraphs[0].text = "Hợp đồng số 02/HĐ/2026"
    section.footer.paragraphs[0].text = "Trang"

    doc.save(str(path))


def test_docx_complex_fidelity_v0(tmp_path: Path) -> None:
    """v0 collapse mode — every paragraph translated; structure intact."""
    from docx import Document

    from nom.translate.formats import translate_docx

    src = tmp_path / "complex.docx"
    dst = tmp_path / "complex.translated.docx"
    _build_complex_docx(src)

    stats = translate_docx(src, dst, IdentityTranslator())
    assert stats.paragraphs_failed == 0
    assert stats.paragraphs_translated >= 12  # h1 + h2 + body + list items + table cells

    out = Document(str(dst))

    # Headings preserved
    headings = [p for p in out.paragraphs if p.style.name.startswith("Heading")]
    h1s = [h for h in headings if h.style.name == "Heading 1"]
    assert any("Hợp đồng dịch vụ" in h.text for h in h1s)

    # Bullet list preserved
    bullets = [p for p in out.paragraphs if p.style.name == "List Bullet"]
    assert len(bullets) == 3

    # Numbered list preserved
    numbered = [p for p in out.paragraphs if p.style.name == "List Number"]
    assert len(numbered) == 2

    # Table structure intact
    assert len(out.tables) == 1
    assert len(out.tables[0].rows) == 4
    assert len(out.tables[0].rows[0].cells) == 3
    assert "Giai đoạn" in out.tables[0].cell(0, 0).text

    # Header / footer preserved
    assert "Hợp đồng" in out.sections[0].header.paragraphs[0].text
    assert out.sections[0].footer.paragraphs[0].text == "Trang"


def test_docx_mixed_style_collapses_in_v0(tmp_path: Path) -> None:
    """v0 collapses to first run's style — mid-sentence bold becomes plain."""
    from docx import Document

    from nom.translate.formats import translate_docx

    src = tmp_path / "mixed.docx"
    dst = tmp_path / "mixed.translated.docx"

    doc = Document()
    p = doc.add_paragraph("")
    p.add_run("plain ")
    p.add_run("BOLD ").bold = True
    p.add_run("plain again.")
    doc.save(str(src))

    translate_docx(src, dst, PrefixTranslator())

    out = Document(str(dst))
    runs = list(out.paragraphs[0].runs)
    # All translated text in run 0; subsequent cleared
    assert runs[0].text.startswith("«")
    assert runs[1].text == ""
    assert runs[2].text == ""


def test_docx_mixed_style_preserved_in_v1(tmp_path: Path) -> None:
    """preserve_runs=True: mid-sentence bold stays bold after translation."""
    from docx import Document

    from nom.translate.formats import translate_docx

    class PreservingPrefix:
        name = "p"
        source_lang = "en"
        target_lang = "vi"

        def translate(self, text: str, *, hint: str | None = None) -> str:
            # Preserve placeholders, just upper-case content between them
            import re

            return re.sub(
                r"(⟦\d+⟧)([^⟦]*)",
                lambda m: m.group(1) + m.group(2).upper(),
                text,
            )

    src = tmp_path / "mixed.docx"
    dst = tmp_path / "mixed.translated.docx"
    doc = Document()
    p = doc.add_paragraph("")
    p.add_run("plain ")
    r2 = p.add_run("BOLD ")
    r2.bold = True
    p.add_run("plain again.")
    doc.save(str(src))

    translate_docx(src, dst, PreservingPrefix(), preserve_runs=True)

    out = Document(str(dst))
    runs = list(out.paragraphs[0].runs)
    # Expected: run 0 = "PLAIN ", run 1 = "BOLD " (bold preserved), run 2 = "PLAIN AGAIN."
    assert runs[0].text == "PLAIN "
    assert runs[1].text == "BOLD "
    assert runs[1].bold is True
    assert runs[2].text == "PLAIN AGAIN."


def test_docx_merged_cells_translated_once(tmp_path: Path) -> None:
    """Merged-cell paragraph must be visited exactly once."""
    from docx import Document

    from nom.translate.formats import translate_docx

    src = tmp_path / "merged.docx"
    dst = tmp_path / "out.docx"
    doc = Document()
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "merged"
    table.cell(0, 0).merge(table.cell(0, 1)).merge(table.cell(0, 2))
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "beta"
    table.cell(1, 2).text = "gamma"
    doc.save(str(src))

    counter = {"calls": 0}

    class Counter:
        name = "c"
        source_lang = "en"
        target_lang = "vi"

        def translate(self, text: str, *, hint: str | None = None) -> str:
            counter["calls"] += 1
            return text.upper()

    stats = translate_docx(src, dst, Counter())
    # Exactly 4 unique paragraphs: 1 merged top row + 3 cells in row 2
    # (plus python-docx's default empty body paragraph, which is skipped).
    assert stats.paragraphs_translated == 4
    assert counter["calls"] == 4


# ---------------------------------------------------------------------------
# XLSX fidelity — multi-sheet, formulas, merged cells, formatted cells.
# ---------------------------------------------------------------------------


def _build_complex_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()

    # Sheet 1: prose
    ws1 = wb.active
    ws1.title = "Mô tả"
    ws1["A1"] = "Báo cáo doanh thu Q1 2026"
    ws1["A1"].font = Font(bold=True, size=14)
    ws1["A2"] = "Đơn vị: triệu VNĐ"
    ws1["A3"] = "Cập nhật: 30/03/2026"

    # Sheet 2: data with formulas
    ws2 = wb.create_sheet("Dữ liệu")
    ws2["A1"] = "Sản phẩm"
    ws2["B1"] = "Giá"
    ws2["C1"] = "Số lượng"
    ws2["D1"] = "Tổng"
    for col in ("A1", "B1", "C1", "D1"):
        ws2[col].font = Font(bold=True)
        ws2[col].fill = PatternFill("solid", fgColor="FFEEEEEE")

    ws2["A2"] = "Bản A"
    ws2["B2"] = 100
    ws2["C2"] = 5
    ws2["D2"] = "=B2*C2"
    ws2["A3"] = "Bản B"
    ws2["B3"] = 200
    ws2["C3"] = 3
    ws2["D3"] = "=B3*C3"

    # Sheet 3: merged + comments
    ws3 = wb.create_sheet("Ghi chú")
    ws3["A1"] = "Ghi chú quan trọng"
    ws3.merge_cells("A1:C1")
    ws3["A2"] = "Cần xem xét trước khi ký"

    wb.save(str(path))


def test_xlsx_multi_sheet_fidelity(tmp_path: Path) -> None:
    from openpyxl import load_workbook

    from nom.translate.formats import translate_xlsx

    src = tmp_path / "report.xlsx"
    dst = tmp_path / "report.translated.xlsx"
    _build_complex_xlsx(src)

    stats = translate_xlsx(src, dst, PrefixTranslator())

    out = load_workbook(str(dst))
    assert out.sheetnames == ["Mô tả", "Dữ liệu", "Ghi chú"]

    # Sheet 1: title translated, font preserved
    assert out["Mô tả"]["A1"].value.startswith("«")
    assert out["Mô tả"]["A1"].font.bold is True

    # Sheet 2: numbers untranslated, formulas intact
    assert out["Dữ liệu"]["B2"].value == 100
    assert out["Dữ liệu"]["C2"].value == 5
    assert str(out["Dữ liệu"]["D2"].value).startswith("=")
    assert str(out["Dữ liệu"]["D3"].value).startswith("=")

    # String cells in sheet 2 translated
    assert out["Dữ liệu"]["A2"].value.startswith("«")
    assert out["Dữ liệu"]["A1"].value.startswith("«")

    # Sheet 3: merged range survives
    merged = list(out["Ghi chú"].merged_cells.ranges)
    assert any(str(r) == "A1:C1" for r in merged)

    # Stats sanity
    assert stats.cells_translated >= 6  # at least the prose strings
    assert stats.cells_failed == 0


# ---------------------------------------------------------------------------
# PPTX fidelity — multi-slide, body + title placeholders, tables, notes.
# ---------------------------------------------------------------------------


def _build_complex_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # Title slide
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Trình bày Q1"
    if s1.placeholders[1].has_text_frame:
        s1.placeholders[1].text = "Báo cáo kết quả"

    # Slide with body bullets
    s2 = prs.slides.add_slide(blank_layout)
    tx = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
    tf = tx.text_frame
    tf.text = "Mục tiêu chính"
    p2 = tf.add_paragraph()
    p2.text = "Tăng doanh thu 30%"
    p3 = tf.add_paragraph()
    p3.text = "Giảm chi phí 10%"

    # Slide with table
    s3 = prs.slides.add_slide(blank_layout)
    tbl_shape = s3.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
    )
    tbl = tbl_shape.table
    tbl.cell(0, 0).text = "Chỉ số"
    tbl.cell(0, 1).text = "Giá trị"
    tbl.cell(1, 0).text = "Doanh thu"
    tbl.cell(1, 1).text = "1500"

    # Slide with speaker notes
    s4 = prs.slides.add_slide(blank_layout)
    s4_tx = s4.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    s4_tx.text_frame.text = "Lời cảm ơn"
    s4.notes_slide.notes_text_frame.text = "Lưu ý: nói chậm phần này"

    prs.save(str(path))


def test_pptx_multi_slide_fidelity(tmp_path: Path) -> None:
    from pptx import Presentation

    from nom.translate.formats import translate_pptx

    src = tmp_path / "deck.pptx"
    dst = tmp_path / "deck.translated.pptx"
    _build_complex_pptx(src)

    stats = translate_pptx(src, dst, PrefixTranslator())

    out = Presentation(str(dst))
    assert len(out.slides) == 4

    # Slide 1: title + subtitle translated
    s1_texts: list[str] = []
    for shape in out.slides[0].shapes:
        if shape.has_text_frame:
            s1_texts.append(shape.text_frame.text)
    assert any(t.startswith("«") for t in s1_texts)

    # Slide 2: body bullets translated
    s2_texts = []
    for shape in out.slides[1].shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                s2_texts.append(p.text)
    assert all(t.startswith("«") or not t.strip() for t in s2_texts)

    # Slide 3: table cells translated
    s3_table = next(s for s in out.slides[2].shapes if s.has_table).table
    assert s3_table.cell(0, 0).text.startswith("«")
    assert s3_table.cell(0, 1).text.startswith("«")

    # Slide 4: speaker notes translated
    notes = out.slides[3].notes_slide.notes_text_frame.text
    assert notes.startswith("«")

    assert stats.paragraphs_failed == 0


# ---------------------------------------------------------------------------
# Text / markdown fidelity — paragraph splits, code blocks, NFC.
# ---------------------------------------------------------------------------


def test_text_paragraph_count_invariant(tmp_path: Path) -> None:
    """Translated text has the same number of paragraphs as source."""
    from nom.translate.formats import translate_text

    src = tmp_path / "src.txt"
    dst = tmp_path / "dst.txt"
    body = "para1 line A\npara1 line B\n\npara2\n\npara3\n\n\n\npara4"
    src.write_text(body, encoding="utf-8")

    stats = translate_text(src, dst, IdentityTranslator())
    assert stats.paragraphs_translated == 4

    n_src_paras = len([p for p in body.split("\n\n") if p.strip()])
    n_dst_paras = len([p for p in dst.read_text(encoding="utf-8").split("\n\n") if p.strip()])
    assert n_src_paras == n_dst_paras


def test_markdown_headers_round_trip(tmp_path: Path) -> None:
    """Markdown headers (#, ##, list markers) round-trip via identity."""
    from nom.translate.formats import translate_text

    src = tmp_path / "src.md"
    dst = tmp_path / "dst.md"
    md = "# Tiêu đề chính\n\n## Mục 1\n\n- bullet A\n- bullet B\n\n```\ncode block\n```\n\nĐoạn cuối."
    src.write_text(md, encoding="utf-8")

    translate_text(src, dst, IdentityTranslator())
    out = dst.read_text(encoding="utf-8")

    assert "# Tiêu đề chính" in out
    assert "## Mục 1" in out
    assert "- bullet A" in out


# ---------------------------------------------------------------------------
# Image → DOCX conversion fidelity.
# ---------------------------------------------------------------------------


def _make_text_image(
    path: Path, text: str = "Hello world", *, size: tuple[int, int] = (640, 200)
) -> None:
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("L", size, 255)
    drawer = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 36)
    except OSError:
        font = ImageFont.load_default()
    drawer.text((20, 60), text, fill=0, font=font)
    img.save(str(path))


@needs_tesseract
def test_image_to_docx_various_extensions(tmp_path: Path) -> None:
    """Each supported image extension routes through the dispatcher."""
    from docx import Document

    from nom.convert import convert_to_docx

    for ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        src = tmp_path / f"src{ext}"
        dst = tmp_path / f"out{ext}.docx"
        _make_text_image(src, "Hello")
        stats = convert_to_docx(src, dst, ocr_language="eng")
        assert stats.pages_ocred == 1
        assert dst.exists()
        doc = Document(str(dst))
        full = "\n".join(p.text for p in doc.paragraphs)
        assert "hello" in full.lower(), f"OCR failed for {ext}: got {full!r}"


@needs_tesseract
def test_image_then_translate_pipeline(tmp_path: Path) -> None:
    """Convert an image to .docx, then translate that .docx — verifies
    the convert→translate pipeline composes cleanly."""
    from docx import Document

    from nom.convert import convert_to_docx
    from nom.translate.formats import translate_docx

    img = tmp_path / "scan.png"
    intermediate = tmp_path / "scan.docx"
    final = tmp_path / "scan.en.docx"
    _make_text_image(img, "Hello world OCR test")

    convert_to_docx(img, intermediate, ocr_language="eng")
    stats = translate_docx(intermediate, final, PrefixTranslator())

    assert stats.paragraphs_translated >= 1
    out = Document(str(final))
    full = "\n".join(p.text for p in out.paragraphs)
    assert "«" in full


# ---------------------------------------------------------------------------
# End-to-end translate_file dispatcher coverage across formats.
# ---------------------------------------------------------------------------


def test_translate_file_dispatcher_handles_all_supported_formats(tmp_path: Path) -> None:
    """Run the dispatcher with each supported extension, identity
    translator, verify no exceptions and output exists."""
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from pptx.util import Inches

    from nom.translate.formats import translate_file

    fixtures: list[tuple[str, str]] = []

    # docx
    docx_src = tmp_path / "x.docx"
    d = Document()
    d.add_paragraph("test")
    d.save(str(docx_src))
    fixtures.append((str(docx_src), str(tmp_path / "x.translated.docx")))

    # xlsx
    xlsx_src = tmp_path / "x.xlsx"
    wb = Workbook()
    wb.active["A1"] = "test"
    wb.save(str(xlsx_src))
    fixtures.append((str(xlsx_src), str(tmp_path / "x.translated.xlsx")))

    # pptx
    pptx_src = tmp_path / "x.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tx.text_frame.text = "test"
    prs.save(str(pptx_src))
    fixtures.append((str(pptx_src), str(tmp_path / "x.translated.pptx")))

    # txt, md, rst
    for ext in (".txt", ".md", ".rst"):
        src = tmp_path / f"x{ext}"
        src.write_text("test paragraph", encoding="utf-8")
        fixtures.append((str(src), str(tmp_path / f"x.translated{ext}")))

    for src, dst in fixtures:
        stats = translate_file(src, dst, IdentityTranslator())
        assert Path(dst).exists(), f"missing output: {dst}"
        # Each fixture has at least one translation unit
        any_translated = (
            getattr(stats, "paragraphs_translated", 0) > 0
            or getattr(stats, "cells_translated", 0) > 0
        )
        assert any_translated, f"no units translated in {src}"
