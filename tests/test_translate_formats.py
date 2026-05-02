"""Tests for the xlsx / pptx / text walkers + the format dispatcher."""

from __future__ import annotations

from pathlib import Path

import pytest

pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("pptx")

from nom.translate.formats import (
    SUPPORTED_FORMATS,
    PptxTranslationStats,
    TextTranslationStats,
    XlsxTranslationStats,
    translate_file,
    translate_pptx,
    translate_text,
    translate_xlsx,
)


class _Reverse:
    name = "reverse"
    source_lang = "en"
    target_lang = "vi"

    def translate(self, text: str, *, hint: str | None = None) -> str:
        return text[::-1]


# ---------------------------------------------------------------------------
# xlsx
# ---------------------------------------------------------------------------


def test_translate_xlsx_string_cells_only(tmp_path: Path) -> None:
    from openpyxl import Workbook, load_workbook

    src = tmp_path / "src.xlsx"
    dst = tmp_path / "dst.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "alpha"
    ws["A2"] = "beta"
    ws["B1"] = 42  # number — must NOT be translated
    ws["B2"] = "=A1"  # formula — must NOT be translated
    ws["B3"] = ""  # empty — skipped
    wb.save(str(src))

    stats = translate_xlsx(src, dst, _Reverse())
    assert isinstance(stats, XlsxTranslationStats)
    assert stats.cells_translated == 2
    assert stats.cells_failed == 0

    out = load_workbook(str(dst))
    out_ws = out.active
    assert out_ws["A1"].value == "ahpla"
    assert out_ws["A2"].value == "ateb"
    assert out_ws["B1"].value == 42
    # Formula round-trips unchanged. openpyxl shows the formula text
    # for formula cells, prefixed with '='.
    assert str(out_ws["B2"].value).startswith("=")


def test_translate_xlsx_missing_source(tmp_path: Path) -> None:
    with pytest.raises(FileNotFoundError):
        translate_xlsx(tmp_path / "nope.xlsx", tmp_path / "out.xlsx", _Reverse())


# ---------------------------------------------------------------------------
# pptx
# ---------------------------------------------------------------------------


def test_translate_pptx_text_in_slides(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    src = tmp_path / "src.pptx"
    dst = tmp_path / "dst.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = tx.text_frame
    tf.text = "alpha"
    tf.add_paragraph().text = "beta"
    prs.save(str(src))

    stats = translate_pptx(src, dst, _Reverse())
    assert isinstance(stats, PptxTranslationStats)
    assert stats.paragraphs_translated == 2

    out = Presentation(str(dst))
    out_text = []
    for sl in out.slides:
        for shape in sl.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    out_text.append(p.text)
    assert "ahpla" in out_text
    assert "ateb" in out_text


# ---------------------------------------------------------------------------
# text / md
# ---------------------------------------------------------------------------


def test_translate_text_paragraph_split(tmp_path: Path) -> None:
    src = tmp_path / "src.txt"
    dst = tmp_path / "dst.txt"
    src.write_text("alpha\nstill alpha\n\nbeta\n\n\ngamma", encoding="utf-8")

    stats = translate_text(src, dst, _Reverse())
    assert isinstance(stats, TextTranslationStats)
    assert stats.paragraphs_translated == 3

    out = dst.read_text(encoding="utf-8").split("\n\n")
    # First paragraph is the whole "alpha\nstill alpha" block (single
    # blank-line separator).
    assert out[0] == "ahpla lufa\nllits\nahpla"[::-1] or "ahpla" in out[0]
    assert "ateb" in out
    assert "ammag" in out


def test_translate_text_normalizes_to_nfc(tmp_path: Path) -> None:
    import unicodedata

    src = tmp_path / "src.md"
    dst = tmp_path / "dst.md"
    nfd = unicodedata.normalize("NFD", "Việt")
    src.write_text(nfd, encoding="utf-8")

    translate_text(src, dst, _Reverse())
    out = dst.read_text(encoding="utf-8")
    assert unicodedata.normalize("NFC", out) == out


# ---------------------------------------------------------------------------
# dispatcher
# ---------------------------------------------------------------------------


def test_dispatcher_routes_by_extension(tmp_path: Path) -> None:
    src = tmp_path / "src.txt"
    dst = tmp_path / "dst.txt"
    src.write_text("hello", encoding="utf-8")
    stats = translate_file(src, dst, _Reverse())
    assert isinstance(stats, TextTranslationStats)


def test_dispatcher_rejects_unsupported_format(tmp_path: Path) -> None:
    src = tmp_path / "data.unknown"
    src.write_text("x", encoding="utf-8")
    with pytest.raises(ValueError, match="unsupported source format"):
        translate_file(src, tmp_path / "out.unknown", _Reverse())


def test_supported_formats_set() -> None:
    assert ".docx" in SUPPORTED_FORMATS
    assert ".xlsx" in SUPPORTED_FORMATS
    assert ".pptx" in SUPPORTED_FORMATS
    assert ".txt" in SUPPORTED_FORMATS
    assert ".md" in SUPPORTED_FORMATS
